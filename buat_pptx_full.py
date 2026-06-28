from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C = type("C",(),{})()
C.PRI = RGBColor(0x0D,0x6E,0xFD)
C.SEC = RGBColor(0x66,0x1A,0xF2)
C.DRK = RGBColor(0x1E,0x29,0x3B)
C.LBG = RGBColor(0xF8,0xFA,0xFC)
C.W = RGBColor(0xFF,0xFF,0xFF)
C.ACC = RGBColor(0x10,0xB9,0x81)
C.WRN = RGBColor(0xFF,0xC1,0x07)
C.RED = RGBColor(0xE7,0x4C,0x3C)
C.GRY = RGBColor(0x95,0xA5,0xA6)
C.GE = RGBColor(0x4F,0x46,0xE5)

def bg(s,c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb=c
def gbg(s):
    b=s.background;f=b.fill;f.gradient()
    f.gradient_stops[0].color.rgb=C.PRI;f.gradient_stops[0].position=0.0
    f.gradient_stops[1].color.rgb=C.GE;f.gradient_stops[1].position=1.0
def rct(s,l,t,w,h,c,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid();sh.fill.fore_color.rgb=c
    if lc:sh.line.color.rgb=lc;sh.line.width=Pt(1)
    else:sh.line.fill.background()
    return sh
def rnd(s,l,t,w,h,c,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid();sh.fill.fore_color.rgb=c
    if lc:sh.line.color.rgb=lc;sh.line.width=Pt(1)
    else:sh.line.fill.background()
    return sh
def ovl(s,l,t,w,h,c,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    sh.fill.solid();sh.fill.fore_color.rgb=c
    if lc:sh.line.color.rgb=lc;sh.line.width=Pt(1.5)
    else:sh.line.fill.background()
    return sh
def di(s,l,t,w,h,c,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.DIAMOND,l,t,w,h)
    sh.fill.solid();sh.fill.fore_color.rgb=c
    if lc:sh.line.color.rgb=lc;sh.line.width=Pt(1.5)
    else:sh.line.fill.background()
    return sh
def ln(s,x1,y1,x2,y2,c,w=Pt(1.5)):
    cn=s.shapes.add_connector(1,x1,y1,x2,y2);cn.line.color.rgb=c;cn.line.width=w
    return cn
def ar(s,x1,y1,x2,y2,c,w=Pt(1.5)):
    cn=s.shapes.add_connector(1,x1,y1,x2,y2);cn.line.color.rgb=c;cn.line.width=w;cn.line.tail=True
    return cn
def tb(s,l,t,w,h,txt,fs=18,c=C.DRK,b=False,a=PP_ALIGN.LEFT,fn='Segoe UI'):
    bx=s.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=c;p.font.bold=b;p.font.name=fn;p.alignment=a
    return bx
def st(shape,txt,fs=12,c=C.DRK,b=False,a=PP_ALIGN.CENTER):
    tf=shape.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=c;p.font.bold=b;p.font.name='Segoe UI';p.alignment=a
    return tf
def bslide(s,tt,bl,st=None):
    gbg(s)
    if st:tb(s,In(0.8),In(0.4),In(11),In(0.6),st,14,RGBColor(0xBB,0xBB,0xBB))
    tb(s,In(0.8),In(0.9),In(11),In(0.8),tt,32,C.W,True)
    rnd(s,In(0.6),In(1.8),In(12),In(5.2),C.W)
    bx=s.shapes.add_textbox(In(1),In(2.1),In(11.3),In(4.6))
    tf=bx.text_frame;tf.word_wrap=True
    for i,b in enumerate(bl):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=b;p.font.size=Pt(18);p.font.color.rgb=C.DRK;p.font.name='Segoe UI';p.space_after=Pt(10)
def In(v):return Inches(v)

# ========== SLIDE 1: COVER ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);gbg(s)
tb(s,In(0.8),In(0.5),In(11),In(0.5),'TUGAS PROJECT PEMOGRAMAN WEB',14,RGBColor(0xBB,0xBB,0xBB))
tb(s,In(0.8),In(2.0),In(11),In(1.2),'CitaStore',60,C.W,True)
tb(s,In(0.8),In(3.2),In(11),In(0.8),'Aplikasi Toko Online Berbasis Web\nDengan Laravel 12 & Bootstrap 5',24,RGBColor(0xDD,0xDD,0xDD))
rct(s,In(0.8),In(4.5),In(3),In(0.06),C.ACC)
tb(s,In(0.8),In(4.9),In(11),In(0.5),'Disusun oleh: [Nama Lengkap] - NIM: [NIM]',16,RGBColor(0xCC,0xCC,0xCC))
tb(s,In(0.8),In(5.4),In(11),In(0.5),'Program Studi [Prodi] - Universitas [Nama Universitas]',14,RGBColor(0xAA,0xAA,0xAA))
tb(s,In(0.8),In(6.2),In(11),In(0.5),'2026',14,RGBColor(0x99,0x99,0x99))
print('Slide 1 done')

# ========== SLIDE 2: DAFTAR ISI ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Daftar Isi',32,C.W,True)
for i,item in enumerate([
    '1.  Latar Belakang','2.  Tujuan dan Manfaat','3.  Teknologi yang Digunakan',
    '4.  Entity Relationship Diagram (ERD)','5.  Struktur Database & Relasi Tabel',
    '6.  Use Case Diagram','7.  Activity Diagram','8.  Class Diagram (UML)',
    '9.  Sequence Diagram','10. Hak Akses & Role Pengguna','11. Fitur untuk Pengguna (User)',
    '12. Fitur untuk Admin','13. Alur Belanja (Flowchart)','14. Arsitektur MVC',
    '15. Tampilan Antarmuka','16. Validasi & Keamanan','17. Kesimpulan dan Saran']):
    tb(s,In(1),In(1.5+i*0.34),In(10),In(0.35),item,17,C.DRK)
print('Slide 2 done')

# ========== SLIDE 3: LATAR BELAKANG ==========
bslide(prs.slides.add_slide(prs.slide_layouts[6]),'Latar Belakang',[
    'Perkembangan e-commerce semakin pesat - belanja online menjadi kebutuhan utama masyarakat.',
    'Masih banyak UMKM dan toko kecil yang belum memiliki platform digital untuk menjual produk.',
    'Kebutuhan akan aplikasi toko online sederhana berbasis web yang mudah dikelola pemilik toko.',
    'Mengimplementasikan materi Pemrograman Web: Laravel (backend) dan Bootstrap (frontend).',
    'Dibangun dengan pendekatan MVC (Model-View-Controller) untuk arsitektur yang rapi dan terstruktur.'])
print('Slide 3 done')

# ========== SLIDE 4: TUJUAN & MANFAAT ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);gbg(s)
tb(s,In(0.8),In(0.4),In(11),In(0.6),'Tujuan & Manfaat',14,RGBColor(0xBB,0xBB,0xBB))
tb(s,In(0.8),In(0.9),In(11),In(0.8),'Tujuan & Manfaat',32,C.W,True)
rnd(s,In(0.6),In(2.0),In(5.8),In(4.8),C.W)
tb(s,In(1),In(2.2),In(5),In(0.5),'Tujuan',24,C.PRI,True)
for i,t in enumerate(['Membangun aplikasi toko online berbasis web yang fungsional dan responsif.','Menerapkan framework Laravel dengan konsep MVC dan routing.','Mengelola data produk, kategori, keranjang, dan pesanan secara terstruktur.','Mengimplementasikan autentikasi pengguna dan role-based access control.']):
    tb(s,In(1.2),In(3.0+i*0.85),In(4.8),In(0.8),f'\\u2022 {t}',15,C.DRK)
rnd(s,In(6.8),In(2.0),In(5.8),In(4.8),C.W)
tb(s,In(7.2),In(2.2),In(5),In(0.5),'Manfaat',24,C.ACC,True)
for i,m in enumerate(['Pemilik toko dapat mengelola produk & pesanan secara digital.','Pelanggan bisa belanja kapan saja dari mana saja.','Memudahkan pencatatan transaksi dan stok barang.','Menjadi portofolio implementasi full-stack web development.']):
    tb(s,In(7.4),In(3.0+i*0.85),In(4.8),In(0.8),f'\\u2022 {m}',15,C.DRK)
print('Slide 4 done')

# ========== SLIDE 5: TEKNOLOGI ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Teknologi yang Digunakan',32,C.W,True)
for i,(nm,dc,cl) in enumerate([('Laravel 12','Framework PHP backend\\nMVC, Eloquent ORM, Blade',C.PRI),('Bootstrap 5','CSS Framework responsif\\nModern & mobile-friendly',C.SEC),('MySQL','Database relasional\\n6 tabel utama',RGBColor(0x00,0x79,0x8F)),('Vite + SCSS','Build tool frontend\\nPreprocessing CSS',C.ACC),('Laravel UI','Auth scaffolding\\nLogin, register, reset password',C.RED),('PHP 8.2+','Backend language\\nType hinting, match expr',RGBColor(0x88,0x99,0xBB))]):
    c=i%3;r=i//3;x=In(0.6+c*4.1);y=In(1.6+r*2.7)
    rnd(s,x,y,In(3.7),In(2.3),C.W);rct(s,x,y,In(3.7),In(0.06),cl)
    tb(s,x+In(0.2),y+In(0.2),In(3.3),In(0.5),nm,20,cl,True)
    tb(s,x+In(0.2),y+In(0.8),In(3.3),In(1.0),dc,13,C.DRK)
print('Slide 5 done')

# ========== SLIDE 6: ERD ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Entity Relationship Diagram (ERD)',32,C.W,True)

def ent(s,x,y,nm,cols,cl=C.PRI):
    w=In(2.35);h=In(0.35+0.26*len(cols))
    rnd(s,x,y,w,h,C.W,C.DRK);rct(s,x,y,w,In(0.35),cl)
    tb(s,x+In(0.05),y+In(0.02),w-In(0.1),In(0.31),nm,11,C.W,True,a=PP_ALIGN.CENTER)
    for i,c in enumerate(cols):
        tb(s,x+In(0.08),y+In(0.38+i*0.26),w-In(0.16),In(0.26),c,8,C.DRK)

cx1=In(0.5);cx2=In(3.6);cx3=In(6.7);cx4=In(10.0);ey=In(1.5)
ent(s,cx1,ey,'users',['PK id','name','email','password','role'],C.PRI)
ent(s,cx2,ey+In(0.5),'carts',['PK id','FK user_id','FK product_id','quantity'],C.ACC)
ent(s,cx3,ey+In(0.0),'categories',['PK id','name'],C.SEC)
ent(s,cx4,ey+In(0.5),'products',['PK id','FK category_id','name','description','price','stock','image'],C.WRN)
oy=In(5.2)
ent(s,cx1,oy,'orders',['PK id','FK user_id','total_price','status','payment_method','shipping_address'],C.RED)
ent(s,cx2,oy+In(0.5),'order_items',['PK id','FK order_id','FK product_id','quantity','price'],RGBColor(0xE6,0x7E,0x22))

# Relationship labels
for lbl,x,y in [('1 ---< N',cx1+In(2.35),ey+In(0.7)),('1 ---< N',cx1+In(2.35),ey+In(3.4)),('1 ---< N',cx2+In(2.35),ey+In(0.7)),('1 ---< N',cx3+In(2.35),ey+In(0.05)),('1 ---< N',In(5.9),oy+In(0.7)),('1 ---< N',In(2.7),oy+In(0.7))]:
    tb(s,x,y,In(1.1),In(0.25),lbl,10,C.DRK,True,a=PP_ALIGN.CENTER)
print('Slide 6 done')

# ========== SLIDE 7: STRUKTUR DB ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Struktur Database & Relasi Tabel',32,C.W,True)
for i,(nm,cols,desc) in enumerate([('users','id, name, email, password, role','Data pengguna & admin'),('categories','id, name','Kategori produk'),('products','id, category_id, name, description,\\nprice, stock, image','Produk dengan gambar'),('carts','id, user_id, product_id, quantity','Keranjang belanja'),('orders','id, user_id, total_price, status,\\npayment_method, shipping_address','Pesanan pelanggan'),('order_items','id, order_id, product_id,\\nquantity, price','Item detail pesanan')]):
    c=i%3;r=i//3;x=In(0.5+c*4.2);y=In(1.5+r*2.8)
    rnd(s,x,y,In(3.8),In(2.5),C.W);rct(s,x,y,In(3.8),In(0.06),C.PRI)
    tb(s,x+In(0.2),y+In(0.15),In(3.4),In(0.4),nm,20,C.PRI,True)
    tb(s,x+In(0.2),y+In(0.65),In(3.4),In(1.0),cols,11,RGBColor(0x55,0x55,0x55),fn='Consolas')
    tb(s,x+In(0.2),y+In(1.85),In(3.4),In(0.4),desc,12,C.DRK)
print('Slide 7 done')

# ========== SLIDE 8: RELASI ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);gbg(s)
tb(s,In(0.8),In(0.4),In(11),In(0.6),'Relasi Antar Tabel',14,RGBColor(0xBB,0xBB,0xBB))
tb(s,In(0.8),In(0.9),In(11),In(0.8),'Hubungan Antar Entitas',32,C.W,True)
for i,(l,rel,r,desc) in enumerate([('Category','1 ---< N','Product','Satu kategori memiliki banyak produk'),('Product','1 ---< N','Cart','Satu produk bisa di-keranjang oleh banyak user'),('Product','1 ---< N','OrderItem','Satu produk muncul di banyak item pesanan'),('User','1 ---< N','Cart','Satu user memiliki banyak item keranjang'),('User','1 ---< N','Order','Satu user dapat membuat banyak pesanan'),('Order','1 ---< N','OrderItem','Satu pesanan terdiri dari banyak item')]):
    c=i%2;r2=i//2;x=In(0.5+c*6.3);y=In(2.0+r2*1.65)
    rnd(s,x,y,In(5.8),In(1.4),C.W)
    tb(s,x+In(0.2),y+In(0.2),In(1.8),In(0.5),l,18,C.PRI,True)
    tb(s,x+In(2.0),y+In(0.2),In(1.8),In(0.5),rel,16,C.ACC,True,a=PP_ALIGN.CENTER)
    tb(s,x+In(3.8),y+In(0.2),In(1.8),In(0.5),r,18,C.SEC,True)
    tb(s,x+In(0.3),y+In(0.8),In(5.2),In(0.4),desc,13,RGBColor(0x55,0x55,0x55))
print('Slide 8 done')

print('\\nPart 1 complete, continuing...')

# ========== SLIDE 9: USE CASE DIAGRAM ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Use Case Diagram',32,C.W,True)

# System boundary
rnd(s,In(3.0),In(1.5),In(10.0),In(5.7),RGBColor(0xE8,0xEC,0xF0),C.DRK)
tb(s,In(3.3),In(1.55),In(3),In(0.3),'CitaStore',12,C.DRK,True)

# Actors
def uc_box(s,x,y,txt,cl):
    b=rnd(s,x,y,In(1.6),In(0.45),cl,cl);st(b,txt,10,C.W,True)

uc_box(s,In(0.4),In(2.8),'Guest',C.GRY)
uc_box(s,In(0.4),In(5.0),'User',C.PRI)
uc_box(s,In(0.4),In(6.2),'Admin',C.WRN)
ar(s,In(1.2),In(5.45),In(1.2),In(3.25),C.DRK)
ar(s,In(1.2),In(6.65),In(1.2),In(5.45),C.DRK)
tb(s,In(1.3),In(4.7),In(1.5),In(0.3),'extends',7,C.GRY)
tb(s,In(1.3),In(5.9),In(1.5),In(0.3),'extends',7,C.GRY)

def uco(s,x,y,lbl,cl):
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,x,y,In(1.3),In(0.4))
    o.fill.solid();o.fill.fore_color.rgb=C.W;o.line.color.rgb=cl;o.line.width=Pt(1.2)
    st(o,lbl,9,cl,True)
def ucln(s,x,y,x2,y2):
    ln(s,x,y,x2,y2,C.GRY,Pt(1))

# Guest use cases
yb=In(1.7)
for lbl,x in [('Lihat Katalog',In(3.3)),('Cari Produk',In(5.3)),('Filter Kategori',In(7.3)),('Registrasi',In(9.3)),('Login',In(11.0))]:
    uco(s,x,yb,lbl,C.PRI);ucln(s,In(2.0),yb+In(0.2),x+In(0.65),yb+In(0.2))
# User use cases
for lbl,x in [('Kelola Keranjang',In(3.3)),('Checkout',In(5.3)),('Lihat Struk',In(7.3))]:
    uco(s,x,yb+In(1.7),lbl,C.ACC);ucln(s,In(2.0),yb+In(1.7)+In(0.2),x+In(0.65),yb+In(1.7)+In(0.2))
# Admin use cases
for lbl,x in [('CRUD Produk',In(3.3)),('CRUD Kategori',In(5.3)),('Kelola Pesanan',In(7.3)),('Update Status',In(9.3))]:
    uco(s,x,yb+In(3.4),lbl,C.WRN);ucln(s,In(2.0),yb+In(3.4)+In(0.2),x+In(0.65),yb+In(3.4)+In(0.2))

# include / extend
tb(s,In(5.0),yb+In(2.1),In(1.2),In(0.2),'<<include>>',7,C.GRY,True,a=PP_ALIGN.CENTER)
ar(s,In(5.3),yb+In(2.2),In(5.3),yb+In(2.5),C.GRY,Pt(1))
# vertical actor line
ln(s,In(2.0),In(2.1),In(2.0),In(5.3),C.GRY,Pt(1))
print('Slide 9 done')

# ========== SLIDE 10: ACTIVITY DIAGRAM ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Activity Diagram - Alur Checkout',32,C.W,True)

# Swimlanes
rct(s,In(0.5),In(1.5),In(5.8),In(5.6),RGBColor(0xE8,0xF0,0xFE))
rct(s,In(6.3),In(1.5),In(6.5),In(5.6),RGBColor(0xE0,0xF5,0xEC))
tb(s,In(0.5),In(1.55),In(5.8),In(0.3),'User',14,C.PRI,True,a=PP_ALIGN.CENTER)
tb(s,In(6.3),In(1.55),In(6.5),In(0.3),'Sistem',14,C.ACC,True,a=PP_ALIGN.CENTER)
ln(s,In(6.3),In(1.5),In(6.3),In(7.1),C.DRK,Pt(1))

def abox(s,x,y,*args):
    if len(args)==4: w,h,txt,cl=args[0],args[1],args[2],args[3]
    else: txt,cl=args;w,h=In(1.8),In(0.4)
    b=rnd(s,x,y,w,h,cl);st(b,txt,10,C.W,True)

ovl(s,In(2.0),In(2.1),In(0.4),In(0.3),C.DRK)
tb(s,In(1.9),In(2.45),In(0.5),In(0.2),chr(8595),14,C.DRK,True,a=PP_ALIGN.CENTER)
abox(s,In(1.2),In(2.7),'Login',C.PRI)
tb(s,In(1.9),In(3.15),In(0.5),In(0.2),chr(8595),14,C.DRK,True,a=PP_ALIGN.CENTER)
abox(s,In(1.2),In(3.4),'Browse & Add to Cart',C.PRI)
tb(s,In(1.9),In(3.85),In(0.5),In(0.2),chr(8595),14,C.DRK,True,a=PP_ALIGN.CENTER)
abox(s,In(1.2),In(4.1),'Klik Checkout',C.PRI)
tb(s,In(4.0),In(4.25),In(0.5),In(0.2),chr(8594),20,C.DRK,True,a=PP_ALIGN.CENTER)

abox(s,In(6.5),In(4.1),In(2.0),In(0.4),'Cek Stok Produk',C.ACC)
tb(s,In(7.4),In(4.55),In(0.5),In(0.2),chr(8595),14,C.DRK,True,a=PP_ALIGN.CENTER)

d=di(s,In(7.0),In(4.8),In(0.9),In(0.55),C.W,C.DRK);st(d,'Stok Cukup?',9,C.DRK,True)

tb(s,In(6.0),In(4.75),In(0.9),In(0.2),chr(8595)+' Tidak',10,C.RED,True,a=PP_ALIGN.CENTER)
abox(s,In(4.8),In(5.7),'Tampilkan Error',C.RED)
tb(s,In(5.5),In(5.25),In(0.5),In(0.4),chr(8592),20,C.RED,True,a=PP_ALIGN.CENTER)

tb(s,In(7.95),In(4.75),In(0.5),In(0.2),chr(8595)+' Ya',10,C.ACC,True,a=PP_ALIGN.CENTER)
abox(s,In(7.2),In(5.6),In(2.2),In(0.4),'Buat Order & Items',C.ACC)
tb(s,In(8.2),In(6.05),In(0.5),In(0.2),chr(8595),14,C.DRK,True,a=PP_ALIGN.CENTER)
abox(s,In(7.2),In(6.3),In(2.2),In(0.4),'Kurangi Stok & Hapus Cart',C.ACC)
tb(s,In(9.0),In(6.45),In(0.5),In(0.2),chr(8594),20,C.DRK,True,a=PP_ALIGN.CENTER)
b2=rnd(s,In(9.6),In(6.3),In(1.6),In(0.4),C.SEC);st(b2,'Tampilkan Struk',10,C.W,True)
ovl(s,In(10.2),In(6.8),In(0.4),In(0.3),C.DRK)
print('Slide 10 done')

# ========== SLIDE 11: CLASS DIAGRAM ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Class Diagram (UML)',32,C.W,True)

def uml(s,x,y,w,nm,attrs,methods,cl):
    hd=rct(s,x,y,w,In(0.32),cl)
    tb(s,x,y+In(0.02),w,In(0.3),nm,10,C.W,True,a=PP_ALIGN.CENTER)
    rct(s,x,y+In(0.32),w,In(0.015),C.DRK)
    ah=In(0.2*len(attrs))
    rct(s,x,y+In(0.335),w,ah,C.W)
    ln(s,x,y+In(0.335),x+w,y+In(0.335),C.DRK,Pt(1))
    for i,a in enumerate(attrs):
        tb(s,x+In(0.04),y+In(0.34+i*0.2),w-In(0.08),In(0.2),a,7,C.DRK)
    my=y+In(0.335)+ah
    if methods:
        ln(s,x,my,x+w,my,C.DRK,Pt(1))
        mh=In(0.2*len(methods))
        rct(s,x,my,w,mh,C.W)
        for i,m in enumerate(methods):
            tb(s,x+In(0.04),my+In(0.01+i*0.2),w-In(0.08),In(0.2),m,7,C.DRK)
    th=In(0.335)+ah+(In(0.2*len(methods)) if methods else In(0))
    # border already covered by components

cx=In(0.3);cy=In(1.5);w=In(2.1)
uml(s,cx,cy,w,'User',['- id: int','- name: string','- email: string','- password: string','- role: string'],['+ carts(): HasMany','+ orders(): HasMany'],C.PRI)
uml(s,cx+In(2.4),cy,w,'Cart',['- id: int','- user_id: int','- product_id: int','- quantity: int'],['+ product(): BelongsTo','+ user(): BelongsTo'],C.ACC)
uml(s,cx+In(4.8),cy,w,'Category',['- id: int','- name: string'],['+ products(): HasMany'],C.SEC)
uml(s,cx+In(7.2),cy,w,'Product',['- id: int','- category_id: int','- name: string','- description: text','- price: int','- stock: int','- image: string'],['+ category(): BelongsTo'],C.WRN)

cy2=In(5.2)
uml(s,cx+In(1.2),cy2,w,'Order',['- id: int','- user_id: int','- total_price: int','- status: string','- payment_method: string','- shipping_address: text'],['+ user(): BelongsTo','+ items(): HasMany'],C.RED)
uml(s,cx+In(3.6),cy2,w,'OrderItem',['- id: int','- order_id: int','- product_id: int','- quantity: int','- price: int'],['+ product(): BelongsTo','+ order(): BelongsTo'],RGBColor(0xE6,0x7E,0x22))

for lbl,x,y in [('1',In(2.35),cy+In(0.7)),('*',In(3.9),cy+In(0.7)),('1',In(4.8),cy+In(1.4)),('*',In(6.6),cy+In(1.4)),('1',In(4.8),cy+In(0.7)),('*',In(6.6),cy+In(0.7)),('1',In(1.2)+In(2.35),cy2+In(0.5)),('*',In(3.6),cy2+In(0.5)),('1',In(3.6)+In(2.35),cy2+In(0.8)),('*',In(6.6),cy2+In(0.8))]:
    tb(s,x,y,In(0.3),In(0.25),lbl,9,C.DRK,True,a=PP_ALIGN.CENTER)
print('Slide 11 done')

# ========== SLIDE 12: SEQUENCE DIAGRAM ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Sequence Diagram - Proses Checkout',32,C.W,True)

lifelines=[('User',In(1.0),C.PRI),('Browser',In(3.5),C.SEC),('CartController',In(6.0),C.ACC),('Order',In(8.5),C.WRN),('OrderItem',In(10.5),C.RED),('Product',In(12.2),RGBColor(0xE6,0x7E,0x22))]

for nm,x,cl in lifelines:
    hd=rnd(s,x-In(0.5),In(1.5),In(1.0),In(0.35),cl);st(hd,nm,9,C.W,True)
    ln(s,x+In(0.0),In(1.85),x+In(0.0),In(6.5),RGBColor(0xCC,0xCC,0xCC),Pt(1))

def sarr(s,x1,x2,y,lbl,cl=C.PRI):
    ar(s,x1,y,x2,y,cl,Pt(1.2))
    mx=(x1+x2)//2
    tb(s,mx-In(0.6),y-In(0.2),In(1.2),In(0.2),lbl,7,cl,True,a=PP_ALIGN.CENTER)
def sbar(s,x,y,h,cl=C.ACC):
    rct(s,x-In(0.03),y,In(0.06),h,cl)

y=In(2.2)
sarr(s,In(1.0),In(3.5),y,'POST /checkout')
sbar(s,In(3.5),y,In(0.5),C.SEC)
y+=In(0.5)
sarr(s,In(3.5),In(6.0),y,'store(request)')
sbar(s,In(6.0),y,In(3.0),C.ACC)
y+=In(0.5)
sarr(s,In(6.0),In(12.2),y,'cekStock()')
y+=In(0.5)
sarr(s,In(12.2),In(6.0),y,'stock tersedia',C.ACC)
y+=In(0.5)

# alt box
ab=rct(s,In(5.6),y,In(5.0),In(1.8),RGBColor(0xFD,0xF2,0xE9))
tb(s,In(5.6),y+In(0.02),In(0.5),In(0.2),'alt',8,C.WRN,True)
ln(s,In(5.6),y+In(0.25),In(10.6),y+In(0.25),C.WRN,Pt(1))
tb(s,In(5.8),y+In(0.3),In(2),In(0.2),'[stok cukup]',8,C.ACC,True)

sarr(s,In(6.0),In(8.5),y+In(0.6),'create()')
sarr(s,In(8.5),In(10.5),y+In(1.0),'createItems()')
sarr(s,In(10.5),In(12.2),y+In(1.4),'decrementStock()')
sarr(s,In(6.0),In(3.5),y+In(1.8),'redirect(receipt)',C.ACC)
y+=In(2.0)
sarr(s,In(3.5),In(1.0),y,'Tampilkan Struk',C.ACC)
print('Slide 12 done')

# ========== SLIDE 13: HAK AKSES ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Hak Akses & Role Pengguna',32,C.W,True)

def role_card(s,x,y,title,cl,items):
    rnd(s,x,y,In(3.8),In(5.2),C.W);rct(s,x,y,In(3.8),In(0.06),cl)
    tb(s,x+In(0.3),y+In(0.25),In(3.2),In(0.5),title,22,cl,True)
    for i,f in enumerate(items):
        tb(s,x+In(0.3),y+In(1.0+i*0.5),In(3.2),In(0.5),chr(8226)+' '+f,14,C.DRK)

role_card(s,In(0.5),In(1.6),'Tamu (Guest)',C.GRY,['Melihat katalog produk','Mencari produk','Filter kategori','Melihat detail produk','Mendaftar akun','Login'])
role_card(s,In(4.7),In(1.6),'User (Terdaftar)',C.PRI,['Semua fitur Guest','Menambahkan ke keranjang','Mengelola keranjang','Checkout & pesan','Melihat struk pesanan','Update profil'])
role_card(s,In(8.9),In(1.6),'Admin',C.WRN,['Semua fitur User','CRUD produk','CRUD kategori','Melihat semua pesanan','Mengubah status pesanan','Manajemen stok'])
print('Slide 13 done')

# ========== SLIDE 14: FITUR PENGGUNA ==========
bslide(prs.slides.add_slide(prs.slide_layouts[6]),'Fitur untuk Pengguna (User)',[
    'Katalog Produk - Grid 3 kolom, pagination 12 produk/halaman.',
    'Pencarian & Filter - Cari nama/deskripsi; filter kategori.',
    'Detail Produk - Info lengkap, stok, harga, produk terkait.',
    'Keranjang Belanja - Tambah/hapus item, update qty + / -.',
    'Checkout - Transfer Bank / COD / QRIS + alamat pengiriman.',
    'Struk Pesanan - Konfirmasi pesanan dengan tombol cetak.'])
print('Slide 14 done')

# ========== SLIDE 15: FITUR ADMIN ==========
bslide(prs.slides.add_slide(prs.slide_layouts[6]),'Fitur untuk Admin',[
    'CRUD Produk - Tambah (upload gambar JPG/PNG), edit, hapus.',
    'Manajemen Kategori - Add/edit/hapus; proteksi hapus jika ada produk.',
    'Manajemen Pesanan - Daftar semua pesanan dengan status & pembayaran.',
    'Detail Pesanan - Item, data pembeli, metode bayar, alamat kirim.',
    'Update Status - Pending -> Success / Cancelled (badge warna).',
    'IsAdmin Middleware - Proteksi halaman admin; non-admin di-redirect.'])
print('Slide 15 done')

# ========== SLIDE 16: ALUR BELANJA ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Alur Belanja (Flowchart)',32,C.W,True)

for i,(num,tt,dc,cl) in enumerate([('1','Browse Produk','Lihat katalog, cari, filter',C.PRI),('2','Detail Produk','Lihat info lengkap & harga',C.SEC),('3','Login / Daftar','Autentikasi pengguna',C.ACC),('4','Tambah ke Keranjang','Item masuk ke keranjang',C.WRN),('5','Checkout','Pilih pembayaran & alamat kirim',C.RED),('6','Struk Pesanan','Konfirmasi & cetak struk',C.PRI)]):
    x=In(0.5+i*2.1);y=In(1.8)
    rnd(s,x,y,In(1.85),In(3.0),C.W);rct(s,x,y,In(1.85),In(0.06),cl)
    ci=ovl(s,x+In(0.65),y+In(0.2),In(0.5),In(0.5),cl);st(ci,num,18,C.W,True)
    tb(s,x+In(0.1),y+In(0.85),In(1.65),In(0.6),tt,14,C.DRK,True,a=PP_ALIGN.CENTER)
    tb(s,x+In(0.1),y+In(1.5),In(1.65),In(0.8),dc,11,RGBColor(0x66,0x66,0x66),a=PP_ALIGN.CENTER)
    if i<5:
        tb(s,x+In(1.75),y+In(1.0),In(0.5),In(0.5),chr(8594),28,RGBColor(0xBB,0xBB,0xBB),True,a=PP_ALIGN.CENTER)

rnd(s,In(0.5),In(5.3),In(12.3),In(1.6),C.W)
tb(s,In(0.8),In(5.4),In(11.7),In(0.4),'Proses Checkout Detail:',16,C.PRI,True)
tb(s,In(0.8),In(5.85),In(11.7),In(0.4),'Buat Order -> Buat OrderItems -> Kurangi stok -> Hapus keranjang -> Redirect ke struk',14,C.DRK)
tb(s,In(0.8),In(6.25),In(11.7),In(0.4),'Validasi: stok diperiksa | Alamat min 10 karakter | Bayar: Transfer/COD/QRIS',12,RGBColor(0x88,0x88,0x88))
print('Slide 16 done')

# ========== SLIDE 17: MVC ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Arsitektur MVC',32,C.W,True)

rnd(s,In(0.5),In(1.6),In(3.8),In(2.3),C.W);rct(s,In(0.5),In(1.6),In(3.8),In(0.06),C.PRI)
tb(s,In(0.8),In(1.85),In(3.2),In(0.4),'Model',22,C.PRI,True)
tb(s,In(0.8),In(2.4),In(3.2),In(1.2),chr(8226)+' Product\n'+chr(8226)+' Category\n'+chr(8226)+' Cart\n'+chr(8226)+' Order\n'+chr(8226)+' OrderItem\n'+chr(8226)+' User',13,C.DRK)

rnd(s,In(4.7),In(1.6),In(3.8),In(2.3),C.W);rct(s,In(4.7),In(1.6),In(3.8),In(0.06),C.ACC)
tb(s,In(5.0),In(1.85),In(3.2),In(0.4),'View (Blade)',22,C.ACC,True)
tb(s,In(5.0),In(2.4),In(3.2),In(1.2),chr(8226)+' layouts/app.blade.php\n'+chr(8226)+' shop/* (4 view)\n'+chr(8226)+' admin/products/* (3)\n'+chr(8226)+' admin/categories/* (3)\n'+chr(8226)+' admin/orders/* (2)\n'+chr(8226)+' auth/* (5 view)',13,C.DRK)

rnd(s,In(8.9),In(1.6),In(3.8),In(2.3),C.W);rct(s,In(8.9),In(1.6),In(3.8),In(0.06),C.WRN)
tb(s,In(9.2),In(1.85),In(3.2),In(0.4),'Controller',22,C.WRN,True)
tb(s,In(9.2),In(2.4),In(3.2),In(1.2),chr(8226)+' ProductController\n'+chr(8226)+' ShopController\n'+chr(8226)+' Admin\\CategoryController\n'+chr(8226)+' Admin\\OrderController\n'+chr(8226)+' HomeController\n'+chr(8226)+' Auth/* (6)',13,C.DRK)

tb(s,In(3.9),In(2.2),In(1.2),In(0.5),chr(8652),30,RGBColor(0xBB,0xBB,0xBB),True,a=PP_ALIGN.CENTER)
tb(s,In(8.1),In(2.2),In(1.2),In(0.5),chr(8652),30,RGBColor(0xBB,0xBB,0xBB),True,a=PP_ALIGN.CENTER)

rnd(s,In(0.5),In(4.3),In(12.3),In(2.8),C.W);rct(s,In(0.5),In(4.3),In(12.3),In(0.06),C.SEC)
tb(s,In(0.8),In(4.5),In(3),In(0.4),'Routing (web.php)',20,C.SEC,True)
tb(s,In(0.8),In(5.0),In(11.5),In(2.0),'Public Routes:\n  GET / -> ShopController@index     |     GET /product/{id} -> ShopController@show\n\nAuthenticated Routes:\n  GET /cart, POST /cart/add/{product}, POST /checkout, GET /receipt/{order}\n\nAdmin Routes (middleware: auth+admin, prefix: /admin):\n  CRUD produk -> ProductController    |    CRUD kategori -> CategoryController\n  Lihat & update status pesanan -> OrderController',12,C.DRK,fn='Consolas')
print('Slide 17 done')

# ========== SLIDE 18: TAMPILAN UI ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Tampilan Antarmuka (UI)',32,C.W,True)

for i,(tt,dc,cl) in enumerate([('Halaman Katalog','Hero carousel, sidebar filter,\\ngrid 3 kolom, search, pagination',C.PRI),('Detail Produk','Breadcrumb, gambar, info produk,\\ntambah ke keranjang, terkait',C.SEC),('Keranjang','Tabel item, +/- qty, hapus,\\nsummary total, form checkout',C.ACC),('Admin Produk','Tabel produk, thumbnail, badge,\\nharga Rupiah, Edit/Hapus',C.WRN),('Admin Pesanan','Tabel pesanan, badge status\\nwarna, detail item & pembeli',C.RED),('Responsive','Navbar sticky, layout mobile,\\nBootstrap 5 grid, print struk',C.PRI)]):
    c=i%3;r=i//3;x=In(0.5+c*4.2);y=In(1.5+r*2.85)
    rnd(s,x,y,In(3.8),In(2.6),C.W);rct(s,x,y,In(3.8),In(0.06),cl)
    tb(s,x+In(0.2),y+In(0.2),In(3.4),In(0.4),tt,18,cl,True)
    tb(s,x+In(0.2),y+In(0.75),In(3.4),In(1.5),dc,12,C.DRK)
print('Slide 18 done')

# ========== SLIDE 19: VALIDASI & KEAMANAN ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,C.LBG)
rct(s,In(0),In(0),In(13.333),In(1.2),C.DRK)
tb(s,In(0.8),In(0.25),In(11),In(0.7),'Validasi & Keamanan',32,C.W,True)

rnd(s,In(0.5),In(1.6),In(6),In(5.5),C.W);rct(s,In(0.5),In(1.6),In(6),In(0.06),C.ACC)
tb(s,In(0.8),In(1.85),In(5.4),In(0.4),'Validasi Input',22,C.ACC,True)
for i,v in enumerate(['Produk: category_id required, name max:255, price integer, image jpeg/png max 2MB','Kategori: name required, unique, max:255','Keranjang: quantity min:1, max: stock + qty','Checkout: method in:transfer,cod,qris; alamat min:10','Registrasi: name, email unique, password min:8, confirmed','Status pesanan: in:Pending,Success,Cancelled']):
    tb(s,In(0.9),In(2.5+i*0.65),In(5.3),In(0.6),chr(8226)+' '+v,12,C.DRK)

rnd(s,In(6.8),In(1.6),In(6),In(5.5),C.W);rct(s,In(6.8),In(1.6),In(6),In(0.06),C.RED)
tb(s,In(7.1),In(1.85),In(5.4),In(0.4),'Fitur Keamanan',22,C.RED,True)
for i,v in enumerate(['CSRF Protection - Semua form @csrf','Admin Middleware - IsAdmin: role=admin','Autentikasi Wajib - Login untuk cart & checkout','Autorisasi Struk - Hanya pemilik pesanan','Cascade Delete - Hapus kategori -> hapus produk','Cek Stok - Cegah pembelian melebihi stok','Filter Gambar - Hanya JPG/PNG max 2MB']):
    tb(s,In(7.2),In(2.5+i*0.65),In(5.3),In(0.6),chr(8226)+' '+v,12,C.DRK)
print('Slide 19 done')

# ========== SLIDE 20: KESIMPULAN ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);gbg(s)
tb(s,In(0.8),In(0.4),In(11),In(0.6),'Penutup',14,RGBColor(0xBB,0xBB,0xBB))
tb(s,In(0.8),In(0.9),In(11),In(0.8),'Kesimpulan & Saran',32,C.W,True)

rnd(s,In(0.6),In(1.9),In(5.8),In(5.0),C.W)
tb(s,In(1),In(2.1),In(5),In(0.5),'Kesimpulan',22,C.PRI,True)
for i,k in enumerate(['Berhasil membangun aplikasi toko online fungsional.','Mengimplementasikan CRUD & role-based access.','Tampilan responsif dengan Bootstrap 5.','Database terstruktur dengan relasi jelas.','Menerapkan MVC Laravel dengan routing lengkap.']):
    tb(s,In(1.2),In(2.8+i*0.65),In(4.8),In(0.6),chr(10003)+' '+k,13,C.DRK)

rnd(s,In(6.8),In(1.9),In(5.8),In(5.0),C.W)
tb(s,In(7.2),In(2.1),In(5),In(0.5),'Saran Pengembangan',22,C.ACC,True)
for i,sg in enumerate(['Sistem kupon diskon & promo.','Notifikasi email checkout & update status.','Upload multiple gambar produk.','Fitur ulasan dan rating produk.','Integrasi payment gateway.']):
    tb(s,In(7.4),In(2.8+i*0.65),In(4.8),In(0.6),chr(8594)+' '+sg,13,C.DRK)
print('Slide 20 done')

# ========== SLIDE 21: TERIMA KASIH ==========
s=prs.slides.add_slide(prs.slide_layouts[6]);gbg(s)
tb(s,In(0.8),In(2.0),In(11),In(1.2),'Terima Kasih',56,C.W,True)
rct(s,In(5.5),In(3.3),In(2.5),In(0.06),C.ACC)
tb(s,In(0.8),In(3.7),In(11),In(0.8),'Ada pertanyaan?',24,RGBColor(0xDD,0xDD,0xDD))
tb(s,In(0.8),In(5.0),In(11),In(0.5),'CitaStore - Aplikasi Toko Online dengan Laravel 12 & Bootstrap 5',16,RGBColor(0xBB,0xBB,0xBB))
tb(s,In(0.8),In(5.5),In(11),In(0.5),'[Nama Lengkap] - [NIM] - [Program Studi] - [Universitas]',14,RGBColor(0x99,0x99,0x99))
print('Slide 21 done')

# ========== SAVE ==========
out=os.path.expanduser('~/Desktop/Presentasi_CitaStore_TokoOnline.pptx')
prs.save(out)
print(f'Saved to: {out}')
print(f'Total slides: {len(prs.slides)}')
